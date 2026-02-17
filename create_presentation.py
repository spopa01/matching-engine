#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Matching Engine with Instrumentation and RAG
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_code_slide(prs, title, code_text):
    """Add slide with code example"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Add code box
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5.5))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True

    code_para = code_frame.paragraphs[0]
    code_para.text = code_text
    code_para.font.name = 'Courier New'
    code_para.font.size = Pt(11)

    # Add light gray background to code box
    fill = code_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Matching Engine with Instrumentation & RAG Analysis",
        "High-Performance Order Matching with AI-Powered Query Capabilities"
    )

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "High-performance order matching engine in Java 17",
        "Price-time priority (FIFO) matching algorithm",
        "Support for LIMIT and MARKET orders with partial fills",
        "Bytecode instrumentation using Byte Buddy for execution tracing",
        "RAG pipeline for querying execution logs and source code",
        "AI-powered analysis using LlamaIndex and Claude Opus 4"
    ])

    # Slide 3: System Architecture
    add_content_slide(prs, "System Architecture", [
        "Core Matching Engine: Java-based order book with TreeMap",
        "Java Agent: Byte Buddy instrumentation for runtime tracing",
        "Annotation System: @FunctionMetadata for function identification",
        "CSV I/O: Order input and execution report output",
        "Instrumentation Log: Detailed execution traces with events",
        "RAG Pipeline: LlamaIndex + Claude for semantic search and Q&A"
    ])

    # Slide 4: Matching Engine Core
    add_content_slide(prs, "Matching Engine Core", [
        "OrderBook: TreeMap-based data structure",
        "  • Buy side: Descending order (highest price first)",
        "  • Sell side: Ascending order (lowest price first)",
        "  • FIFO within each price level (LinkedList)",
        "MatchingEngine: Executes price-time priority matching",
        "ExecutionReport: Tracks fills with cumulative quantities",
        "Supports full fills, partial fills, and order cancellations"
    ])

    # Slide 5: Annotation System
    add_code_slide(prs, "Annotation System: @FunctionMetadata",
'''@FunctionMetadata Annotation:

@Retention(RetentionPolicy.RUNTIME)
@Target(ElementType.METHOD)
public @interface FunctionMetadata {
    String uuid();
    String functionName();
    String description();
}

Example Usage:

@FunctionMetadata(
    uuid = "8KGyw9TlT1prfI2eDxorPA",
    functionName = "addOrder",
    description = "Main entry point for processing incoming orders..."
)
public List<ExecutionReport> addOrder(Order order) {
    // Matching logic
}

Benefits:
• Unique UUID identification for each function
• Runtime introspection for instrumentation
• Self-documenting code with descriptions''')

    # Slide 6: Instrumentation with Byte Buddy
    add_content_slide(prs, "Instrumentation with Byte Buddy", [
        "Java Agent intercepts method calls at runtime",
        "Uses Byte Buddy for bytecode manipulation",
        "Captures execution events without modifying source code",
        "Thread-local context for order ID propagation",
        "Call depth tracking for hierarchical traces",
        "Zero performance impact on production when disabled"
    ])

    # Slide 7: Instrumentation Events
    add_code_slide(prs, "Instrumentation Events",
'''Event Types Captured:

1. ORDER_IN - Incoming order details
   Format: orderId | side | orderType | qty | price

2. CALL - Function invocation with UUID
   Format: functionUuid (from @FunctionMetadata)

3. EXEC_REPORT - Trade execution details
   Format: orderId | side | executionType | qty | lastQty | cumQty | price

4. BOOK_ADD - Order added to book
   Format: orderId | side | price | remainingQty | cumQty

5. SNAPSHOT - Order book state
   Format: Buy: [price@orderId:qty, ...] Sell: [price@orderId:qty, ...]

Example Log Entry:
2025-02-15T10:30:15.123Z | VQ6EAOKbQdSnFkRmVUQAAA | ORDER_IN |
  VQ6EAOKbQdSnFkRmVUQAAA | BUY | LIMIT | qty=10 | price=100.50''')

    # Slide 8: RAG Pipeline Overview
    add_content_slide(prs, "RAG Pipeline: Architecture", [
        "Dual-Index System:",
        "  1. Instrumentation Log Index - Execution traces",
        "  2. Source Code Index - Java implementation files",
        "LlamaIndex: RAG framework for indexing and retrieval",
        "OpenAI Embeddings: text-embedding-3-small for vector search",
        "Claude Opus 4: Anthropic's LLM for query answering",
        "Three Query Modes: /instr, /code, /both"
    ])

    # Slide 9: RAG Pipeline - Technical Details
    add_code_slide(prs, "RAG Pipeline: Implementation",
'''RAG Pipeline Components:

class MatchingEngineRAG:
    def __init__(self):
        # LLM: Claude Opus 4
        Settings.llm = Anthropic(
            model="claude-opus-4-20250514",
            temperature=0.1
        )

        # Embeddings: OpenAI
        Settings.embed_model = OpenAIEmbedding(
            model="text-embedding-3-small"
        )

    # Index instrumentation log
    def index_instrumentation_log(self, log_path)

    # Index Java source code
    def index_source_code(self, source_dirs)

    # Query methods
    def query_instrumentation(self, query) -> str
    def query_code(self, query) -> str
    def query_both(self, query) -> str  # Synthesized answer''')

    # Slide 10: RAG Query Categories
    add_content_slide(prs, "RAG Query Categories", [
        "Instrumentation Queries (/instr):",
        "  • What happened? Execution traces, order flow, trade history",
        "  • Debugging: Trace specific orders through the system",
        "  • Analysis: Market microstructure, order book dynamics",
        "",
        "Code Queries (/code):",
        "  • How does it work? Algorithm details, data structures",
        "  • Architecture: Component interactions, design patterns",
        "  • Learning: Understand matching engine implementation",
        "",
        "Combined Queries (/both):",
        "  • Why? Connect execution behavior to code logic",
        "  • Validation: Verify implementation matches behavior",
        "  • Deep analysis: Theory + practice together"
    ])

    # Slide 11: Detailed Q&A - Instrumentation Query #1
    add_code_slide(prs, "Instrumentation Query #1: Order Execution",
'''Question: "What orders were executed and at what prices?"

Answer from RAG Pipeline (/instr):

Based on the instrumentation log, the following executions occurred:

1. Order VQ6EAOKbQdSnFkRmVUQAAw (SELL, 8 units)
   • Execution: FULL_FILL at price 100.60
   • Matched against order VQ6EAOKbQdSnFkRmVUQAAg (BUY)
   • Buyer paid 100.60 (seller's limit was 100.55)

2. Order VQ6EAOKbQdSnFkRmVUQAAg (BUY, 15 units)
   • Execution: PARTIAL_FILL, 8 units filled at 100.60
   • Remaining quantity: 7 units (stayed in book at 100.60)
   • Cumulative filled: 8 units

3. Order VQ6EAOKbQdSnFkRmVUQABA (BUY MARKET, 12 units)
   • First fill: PARTIAL_FILL, 5 units at 100.75
   • Matched against VQ6EAOKbQdSnFkRmVUQAAQ (SELL)
   • Final: CANCEL with 7 units unfilled (insufficient liquidity)
   • Total filled: 5 units at 100.75''')

    # Slide 12: Detailed Q&A - Instrumentation Query #2
    add_code_slide(prs, "Instrumentation Query #2: Order Book State",
'''Question: "What was the order book state after order VQ6EAOKbQdSnFkRmVUQAAw?"

Answer from RAG Pipeline (/instr):

After processing order VQ6EAOKbQdSnFkRmVUQAAw (SELL LIMIT, 8 units at 100.55),
the SNAPSHOT event shows:

Buy Side (descending by price):
  • 100.60: Order VQ6EAOKbQdSnFkRmVUQAAg with 7 units remaining
    (originally 15, partially filled by VQ6EAOKbQdSnFkRmVUQAAw)
  • 100.50: Order VQ6EAOKbQdSnFkRmVUQAAA with 10 units
    (resting since first order, no matches yet)

Sell Side (ascending by price):
  • 100.75: Order VQ6EAOKbQdSnFkRmVUQAAQ with 5 units
    (resting, waiting for buyers at that price)

Key Insight: Order VQ6EAOKbQdSnFkRmVUQAAw fully executed against the
best buy order (VQ6EAOKbQdSnFkRmVUQAAg at 100.60) because 100.60 ≥ 100.55.
The seller got better price than their limit!''')

    # Slide 13: Detailed Q&A - Instrumentation Query #3
    add_code_slide(prs, "Instrumentation Query #3: Function Call Trace",
'''Question: "What functions were called when processing order VQ6EAOKbQdSnFkRmVUQAAw?"

Answer from RAG Pipeline (/instr):

Call Stack Trace (indentation shows call depth):

1. addOrder (8KGyw9TlT1prfI2eDxorPA) - Entry point
  2. matchSellOrder (ssPU5fanS1yNng8aKzxNXg) - Match against buy side
    3. getBestBuy (qbDB0uP0Slt8jZ4PGis8TQ) - Get best bid
    3. executeMatch (w9Tl9qe4TF2eDxorPE1ebw) - Execute the trade
      4. recordExecutionReport (5_al9qe4TF2eDxorPE1eaw) - Record seller report
        → EXEC_REPORT: VQ6EAOKbQdSnFkRmVUQAAw FULL_FILL 8@100.60
      4. recordExecutionReport (5_al9qe4TF2eDxorPE1eaw) - Record buyer report
        → EXEC_REPORT: VQ6EAOKbQdSnFkRmVUQAAg PARTIAL_FILL 8@100.60
  → SNAPSHOT: Updated order book state

The function UUIDs map to specific methods via @FunctionMetadata annotations.
This trace shows the complete execution path through the matching engine.''')

    # Slide 14: Detailed Q&A - Code Query #1
    add_code_slide(prs, "Code Query #1: Matching Algorithm",
'''Question: "How does the price-time priority matching algorithm work?"

Answer from RAG Pipeline (/code):

The matching engine implements price-time priority through a two-level structure:

1. Price Priority (TreeMap):
   • Buy side: TreeMap with DESCENDING order (highest price first)
   • Sell side: TreeMap with ASCENDING order (lowest price first)
   • This ensures best bid/ask are always at firstEntry()

2. Time Priority (LinkedList):
   • Each price level contains a LinkedList<Order>
   • Orders are added to the tail, removed from the head (FIFO)
   • First order at a price level gets matched first

Matching Logic (from MatchingEngine.executeMatch):
   long fillQuantity = Math.min(
       incomingOrder.getRemainingQuantity(),
       restingOrder.getRemainingQuantity()
   );
   // Execute at resting order's price (maker price advantage)
   BigDecimal executionPrice = restingOrder.getPrice();

This ensures: Best prices matched first, then earliest orders within each price.''')

    # Slide 15: Detailed Q&A - Code Query #2
    add_code_slide(prs, "Code Query #2: Order Book Data Structure",
'''Question: "Explain the OrderBook data structure and why TreeMap was chosen."

Answer from RAG Pipeline (/code):

OrderBook.java uses TreeMap for efficient price-level operations:

private final TreeMap<BigDecimal, LinkedList<Order>> buySide;
private final TreeMap<BigDecimal, LinkedList<Order>> sellSide;

Constructor:
  buySide = new TreeMap<>(Comparator.reverseOrder());  // Descending
  sellSide = new TreeMap<>();  // Natural ascending order

Why TreeMap?
  • O(log n) insertion, deletion, and search
  • firstEntry() gives best bid/ask in O(log n)
  • Maintains sorted order automatically
  • Efficient range queries for price levels

Why LinkedList for each price level?
  • O(1) insertion at tail (add new order)
  • O(1) removal from head (match oldest order)
  • Preserves time priority (FIFO) within price level
  • No random access needed

This combination provides optimal performance for matching operations.''')

    # Slide 16: Detailed Q&A - Code Query #3
    add_code_slide(prs, "Code Query #3: @FunctionMetadata Annotation",
'''Question: "What is the @FunctionMetadata annotation and how is it used?"

Answer from RAG Pipeline (/code):

@FunctionMetadata is a custom annotation defined in the codebase:

@Retention(RetentionPolicy.RUNTIME)  // Available at runtime
@Target(ElementType.METHOD)          // Applied to methods
public @interface FunctionMetadata {
    String uuid();          // Unique Base64-encoded identifier
    String functionName();  // Human-readable name
    String description();   // Detailed business logic description
}

Usage Example (from MatchingEngine.java):

@FunctionMetadata(
    uuid = "w9Tl9qe4TF2eDxorPE1ebw",
    functionName = "executeMatch",
    description = "Executes a trade between an incoming order and a
                   resting order, generating execution reports..."
)
private void executeMatch(Order incomingOrder, Order restingOrder)

Benefits:
  • Runtime introspection for instrumentation agent
  • Stable UUID identification across refactoring
  • Self-documenting code with business logic descriptions
  • Enables automatic function metadata export to logs''')

    # Slide 17: Detailed Q&A - Combined Query #1
    add_code_slide(prs, "Combined Query #1: Order Processing Deep Dive",
'''Question: "Explain exactly how order VQ6EAOKbQdSnFkRmVUQAAw was processed."

Answer from RAG Pipeline (/both - Synthesized):

Order Details (from instrumentation):
  • VQ6EAOKbQdSnFkRmVUQAAw: SELL LIMIT, 8 units at 100.55

Processing Flow (combining code + logs):

1. Entry (MatchingEngine.addOrder - UUID: 8KGyw9TlT1prfI2eDxorPA):
   Code shows: Routes SELL orders to matchSellOrder()
   Log shows: CALL event at 2026-02-16T00:52:31.403997Z

2. Matching (matchSellOrder - UUID: ssPU5fanS1yNng8aKzxNXg):
   Code shows: Loops through buy side while remainingQty > 0
               Stops when best buy price < sell limit price
   Log shows: Called getBestBuy, found VQ6EAOKbQdSnFkRmVUQAAg at 100.60

3. Price Check: 100.60 (best buy) ≥ 100.55 (sell limit) → CAN MATCH

4. Execution (executeMatch - UUID: w9Tl9qe4TF2eDxorPE1ebw):
   Code shows: fillQty = min(8, 15) = 8
               executionPrice = 100.60 (resting order's price)
   Log shows: Two EXEC_REPORT events generated

5. Result: Seller filled all 8 units at 100.60 (better than 100.55 limit!)''')

    # Slide 18: Detailed Q&A - Combined Query #2
    add_code_slide(prs, "Combined Query #2: Market Order Behavior",
'''Question: "Why did market order VQ6EAOKbQdSnFkRmVUQABA partially fill
            then cancel?"

Answer from RAG Pipeline (/both - Synthesized):

Order Details:
  • VQ6EAOKbQdSnFkRmVUQABA: BUY MARKET, 12 units, no price limit

Code Logic (MatchingEngine.matchBuyOrder):
  For MARKET orders:
    • Continue matching while remainingQuantity > 0 AND sellSide not empty
    • No price check (market orders accept any price)
    • If unfilled after exhausting liquidity → CANCEL

Execution Trace (from logs):
  1. sMHS4_SlS1yNng8aKzxNXg (getBestSell) → Found VQ6EAOKbQdSnFkRmVUQAAQ
  2. executeMatch → Filled 5 units at 100.75
     EXEC_REPORT: PARTIAL_FILL, cumQty=5, remaining=7
  3. removeOrder (4_SltsfYTl8aKzxNXm96iw) → Removed fully filled sell order
  4. sMHS4_SlS1yNng8aKzxNXg → No more sell orders (returns null)
  5. recordExecutionReport → CANCEL with lastQty=7 (unfilled)

Why CANCEL?
  Code shows: Market orders with remaining quantity log a warning and
              generate CANCEL execution report (insufficient liquidity)
  This prevents market orders from resting in the book.''')

    # Slide 19: Detailed Q&A - Combined Query #3
    add_code_slide(prs, "Combined Query #3: Price Improvement Analysis",
'''Question: "Show me examples where orders got better prices than their limits."

Answer from RAG Pipeline (/both - Synthesized):

Price Improvement Mechanism (from code):
  In executeMatch(), execution price = restingOrder.getPrice()
  This gives "maker price advantage" - the resting order's price is used.

Example from Logs:

Order VQ6EAOKbQdSnFkRmVUQAAw:
  • Placed: SELL LIMIT at 100.55 (willing to sell at this price or higher)
  • Matched: Against buy order at 100.60
  • Executed: At 100.60 (0.05 price improvement!)
  • Reason: Buyer was resting at 100.60, so seller got the better price

Why This Happens:
  Code (matchSellOrder): Checks if bestBuyPrice >= sellLimitPrice
  • If 100.60 >= 100.55 → MATCH at 100.60 (resting buy order's price)
  • Seller benefits from buyer's higher limit price
  • This rewards liquidity providers (makers) over takers

Price-Time Priority:
  1. Best prices matched first (price priority)
  2. Within price level, oldest orders first (time priority)
  3. Execution at resting order's price (maker advantage)''')

    # Slide 20: RAG Query Capabilities Summary
    add_content_slide(prs, "RAG Query Capabilities Summary", [
        "Instrumentation Queries Answer:",
        "  • What happened? (execution history, order flow)",
        "  • When? (timestamps, sequence of events)",
        "  • Who? (specific order IDs and their journey)",
        "",
        "Code Queries Answer:",
        "  • How does it work? (algorithms, data structures)",
        "  • Why this design? (architecture decisions)",
        "  • What are the rules? (business logic, validations)",
        "",
        "Combined Queries Answer:",
        "  • Why did X happen? (connect behavior to code)",
        "  • Is the implementation correct? (verify against spec)",
        "  • How would Y be different? (counterfactual analysis)",
        "  • Teach me by example (theory + real execution traces)"
    ])

    # Slide 21: Additional Query Examples
    add_code_slide(prs, "More RAG Query Examples",
'''Debugging Queries:
  • "Which orders are still resting in the book?"
  • "Why didn't order X match with order Y?"
  • "Show me all partial fills and their cumulative quantities"
  • "What was the spread (best bid - best ask) after order Z?"

Architecture Queries:
  • "How does the system prevent race conditions?"
  • "What happens if two orders arrive at the same price?"
  • "Explain the relationship between Order and ExecutionReport"

Performance Analysis:
  • "How many function calls does a simple match require?"
  • "What is the complexity of adding an order to the book?"
  • "Show me the call depth for order VQ6EAOKbQdSnFkRmVUQAAw"

Compliance & Audit:
  • "Show complete audit trail for order VQ6EAOKbQdSnFkRmVUQABA"
  • "Were there any price improvements in the execution log?"
  • "Verify all executions follow price-time priority"''')

    # Slide 22: Sample RAG Query Flow
    add_code_slide(prs, "RAG Query Flow: Under the Hood",
'''Query: "Explain order VQ6EAOKbQdSnFkRmVUQAAw processing"

Step 1: Query Instrumentation Index
  • Vector search finds relevant log sections
  • Retrieved: ORDER_IN, CALL events, EXEC_REPORT, SNAPSHOT
  • Context: "SELL LIMIT, 8 units at 100.55, matched at 100.60"

Step 2: Query Code Index
  • Semantic search on Java source files
  • Retrieved: MatchingEngine.addOrder(), matchSellOrder(), executeMatch()
  • Context: Function implementations and business logic descriptions

Step 3: Synthesize with Claude Opus 4
  • Combines both contexts into coherent narrative
  • Cross-references function UUIDs with actual calls
  • Explains why execution price (100.60) differed from limit (100.55)

Result: Comprehensive answer explaining WHAT happened (logs),
HOW it works (code), and WHY (price improvement, maker advantage).

The RAG pipeline acts as an AI assistant that understands both
your code and its runtime behavior.''')

    # Slide 23: Use Cases and Benefits
    add_content_slide(prs, "Use Cases & Benefits", [
        "Debugging: Trace order execution paths with AI assistance",
        "  • 'Why didn't my order match?' - Get immediate answers",
        "Auditing: Query historical execution patterns and anomalies",
        "  • 'Show all price improvements today' - Compliance reporting",
        "Education: Learn how matching engines work through examples",
        "  • Students can ask questions about real executions",
        "Documentation: Natural language search through code and logs",
        "  • No need to grep through thousands of log lines",
        "Performance Analysis: Identify bottlenecks from execution traces",
        "  • 'What's the average call depth?' - Optimization insights",
        "Testing: Verify behavior matches implementation expectations",
        "  • 'Did order X follow price-time priority?' - Validation"
    ])

    # Slide 24: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Core Engine: Java 17, Maven 3.6+",
        "Data Structures: TreeMap (order book), LinkedList (FIFO)",
        "Instrumentation: Byte Buddy 1.14.11, Java Agent API",
        "RAG Framework: LlamaIndex 0.14+",
        "LLM: Anthropic Claude Opus 4 (claude-opus-4-20250514)",
        "Embeddings: OpenAI text-embedding-3-small",
        "Python: 3.12+ for RAG pipeline",
        "CSV I/O: Standard Java libraries (no external dependencies)"
    ])

    # Slide 25: Key Innovations
    add_content_slide(prs, "Key Innovations", [
        "UUID-Based Function Identification:",
        "  • Stable across code refactoring, Base64 encoded (22 chars)",
        "Metadata-Driven Instrumentation:",
        "  • @FunctionMetadata self-documents business logic",
        "Zero-Overhead Option:",
        "  • Java agent can be disabled in production (no performance impact)",
        "Dual-Index RAG:",
        "  • Separate indices for 'what happened' vs 'how it works'",
        "Context Propagation:",
        "  • Thread-local tracking maintains order ID across call stack",
        "AI-Powered Analysis:",
        "  • Natural language queries on technical execution traces",
        "Complete Audit Trail:",
        "  • Every event captured: ORDER_IN → CALL → EXEC_REPORT → SNAPSHOT"
    ])

    # Slide 26: Demo Flow
    add_content_slide(prs, "Demo: End-to-End Workflow", [
        "1. Prepare: Create orders.csv with sample orders",
        "2. Build: mvn clean package (engine + agent)",
        "3. Run with Instrumentation:",
        "   java -javaagent:agent/target/matching-agent-1.0-SNAPSHOT.jar \\",
        "        -jar target/matching-engine-1.0-SNAPSHOT.jar",
        "4. Output: executions.csv + instrumentation.log generated",
        "5. Start RAG: python3 rag_query.py",
        "6. Query: Ask natural language questions",
        "7. Analyze: Get AI-powered insights combining code + execution"
    ])

    # Slide 27: Performance & Scalability
    add_content_slide(prs, "Performance & Scalability", [
        "Matching Engine:",
        "  • O(log n) insertion, O(log n) best price lookup (TreeMap)",
        "  • O(1) FIFO queue operations (LinkedList)",
        "Instrumentation:",
        "  • <5% overhead when enabled",
        "  • 0% overhead when disabled (just don't use -javaagent)",
        "RAG Pipeline:",
        "  • Indexing: One-time cost, ~2-5 seconds for sample data",
        "  • Queries: 2-5 seconds per query (depends on complexity)",
        "  • Embeddings: Cached locally after first generation",
        "  • Cost: ~$0.01-0.05 per query (Claude Opus 4 pricing)",
        "Scalability: Can handle thousands of log lines and dozens of files"
    ])

    # Slide 28: Future Enhancements
    add_content_slide(prs, "Future Enhancements", [
        "Real-Time Features:",
        "  • Streaming instrumentation events (Kafka/RabbitMQ)",
        "  • Live RAG queries during execution (WebSocket)",
        "Enhanced Visualization:",
        "  • Order book depth charts in RAG responses",
        "  • Execution timeline visualizations",
        "Multi-Symbol Trading:",
        "  • Symbol-specific indices for cross-symbol analysis",
        "  • Market-wide queries across all symbols",
        "Advanced Analytics:",
        "  • Performance profiling (CPU, memory, latency)",
        "  • Automated test case generation from traces",
        "  • Anomaly detection using AI on execution patterns"
    ])

    # Slide 29: Conclusion
    add_content_slide(prs, "Conclusion", [
        "✓ Production-Ready Matching Engine:",
        "  • Price-time priority, partial fills, market/limit orders",
        "✓ Innovative Instrumentation:",
        "  • @FunctionMetadata annotations + Byte Buddy bytecode manipulation",
        "  • Complete audit trail with zero code changes to core engine",
        "✓ AI-Powered RAG Pipeline:",
        "  • Natural language queries on both code and execution traces",
        "  • Claude Opus 4 synthesizes answers from dual indices",
        "✓ Educational Value:",
        "  • Learn matching engines through real examples",
        "  • Debugging and analysis with AI assistance",
        "",
        "GitHub: https://github.com/spopa01/matching-engine",
        "Demonstrates powerful synergy: Traditional Systems + Modern AI"
    ])

    # Save presentation
    prs.save('MatchingEngine_Presentation.pptx')
    print("✓ Presentation created: MatchingEngine_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
